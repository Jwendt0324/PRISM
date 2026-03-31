#!/usr/bin/env python3
"""
Claude PRISM — Onboarding Presentation Deck Builder
Generates a professional .pptx presentation for team/client onboarding.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──
BLACK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_BLUE = RGBColor(0x16, 0x21, 0x3E)
ACCENT_BLUE = RGBColor(0x0F, 0x6F, 0xFF)
ACCENT_TEAL = RGBColor(0x00, 0xD2, 0xD3)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
DARK_GRAY = RGBColor(0x44, 0x44, 0x44)
GREEN = RGBColor(0x00, 0xC9, 0x7B)
ORANGE = RGBColor(0xFF, 0x97, 0x40)
RED = RGBColor(0xFF, 0x61, 0x61)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=BLACK):
    """Set solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a colored rectangle."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, size=16,
                    color=WHITE, spacing=Pt(8)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return tf


# ════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BLACK)

# Accent bar at top
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

# Title
add_text(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.5),
         "THE CLAUDE PRISM", size=54, color=WHITE, bold=True)

# Subtitle
add_text(slide, Inches(1), Inches(3.2), Inches(10), Inches(0.8),
         "A Plug-and-Play Operating System for Claude Code", size=28,
         color=ACCENT_TEAL)

# Description
add_text(slide, Inches(1), Inches(4.4), Inches(8), Inches(1),
         "SOPs, skills, frameworks, and automation — personalized to you\nthrough an identity scan of your environment.",
         size=18, color=MEDIUM_GRAY)

# Bottom bar
add_shape(slide, Inches(0), Inches(6.9), W, Inches(0.08), ACCENT_TEAL)
add_text(slide, Inches(1), Inches(7.0), Inches(6), Inches(0.4),
         "[Your Agency Name]  ×  [Methodology Partner]", size=12, color=MEDIUM_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 2: The Problem
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "THE PROBLEM", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
         "Claude is powerful. But every session starts from zero.", size=36,
         color=WHITE, bold=True)

problems = [
    "Every new conversation, you re-explain who you are and what you do",
    "No memory of your SOPs, frameworks, or methodology",
    "No awareness of your clients, team, or priorities",
    "No quality standards — output varies wildly between sessions",
    "You can't share what works with your team or clients",
]

for i, problem in enumerate(problems):
    y = Inches(2.8) + Inches(i * 0.75)
    # Red X
    add_text(slide, Inches(1), y, Inches(0.5), Inches(0.5),
             "✗", size=22, color=RED, bold=True)
    add_text(slide, Inches(1.5), y, Inches(10), Inches(0.5),
             problem, size=20, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 3: The Solution
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "THE SOLUTION", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1),
         "The PRISM gives Claude a brain before the conversation starts.",
         size=34, color=WHITE, bold=True)

solutions = [
    "Loads your personal context automatically — who you are, what you do",
    "Routes tasks to the right SOP or skill — no guessing",
    "Enforces quality standards — 18-step QA, canon compliance",
    "Compounds over time — every session makes it smarter",
    "Shareable — your team runs the same system, customized to them",
]

for i, solution in enumerate(solutions):
    y = Inches(2.8) + Inches(i * 0.75)
    add_text(slide, Inches(1), y, Inches(0.5), Inches(0.5),
             "✓", size=22, color=GREEN, bold=True)
    add_text(slide, Inches(1.5), y, Inches(10), Inches(0.5),
             solution, size=20, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 4: Two-Layer Architecture
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "HOW IT WORKS", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Two Layers: Shared Core + Personal Identity", size=36,
         color=WHITE, bold=True)

# Left box — Shared Core
box1 = add_shape(slide, Inches(1), Inches(2.5), Inches(5), Inches(4.2), DARK_BLUE)
add_text(slide, Inches(1.4), Inches(2.7), Inches(4), Inches(0.5),
         "SHARED CORE", size=20, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1.4), Inches(3.2), Inches(4.2), Inches(0.4),
         "Same for everyone — synced via Git", size=14, color=MEDIUM_GRAY)

shared_items = [
    "12 Skills — action-oriented Claude prompts",
    "25 SOPs — step-by-step procedures",
    "15 Canon Files — source-of-truth frameworks",
    "Team-Ops Templates — directory, roles, matrix",
    "Scripts & Hooks — automation infrastructure",
]
add_bullet_list(slide, Inches(1.4), Inches(3.8), Inches(4.2), Inches(3),
                shared_items, size=15, color=LIGHT_GRAY, spacing=Pt(10))

# Right box — Personal Identity
box2 = add_shape(slide, Inches(7), Inches(2.5), Inches(5), Inches(4.2), DARK_BLUE)
add_text(slide, Inches(7.4), Inches(2.7), Inches(4), Inches(0.5),
         "PERSONAL IDENTITY", size=20, color=ACCENT_TEAL, bold=True)
add_text(slide, Inches(7.4), Inches(3.2), Inches(4.2), Inches(0.4),
         "Unique to you — stays on your machine", size=14, color=MEDIUM_GRAY)

personal_items = [
    "CONTEXT.md — who you are, what you do",
    "Memory Bank — clients, team, priorities",
    "Communication Style — tone, sign-off, patterns",
    "Active Projects — what's happening now",
    "Auto-generated by Identity Scan",
]
add_bullet_list(slide, Inches(7.4), Inches(3.8), Inches(4.2), Inches(3),
                personal_items, size=15, color=LIGHT_GRAY, spacing=Pt(10))


# ════════════════════════════════════════════════════
# SLIDE 5: Identity Scan Overview
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "THE IDENTITY SCAN", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Claude reads your world and builds your context automatically.",
         size=34, color=WHITE, bold=True)

# Source cards
sources = [
    ("💻", "Computer", "Projects, documents,\ncode repos, configs", ACCENT_BLUE),
    ("📧", "Gmail", "Contacts, threads,\ncommunication style", GREEN),
    ("📅", "Calendar", "Meetings, patterns,\nkey relationships", ORANGE),
    ("📁", "Google Drive", "Shared docs, client\nfolders, assets", ACCENT_TEAL),
    ("🔧", "GitHub", "Repos, contributions,\ncollaborators", RGBColor(0xAB, 0x6E, 0xFF)),
    ("📝", "Notion", "Workspaces, databases,\nproject notes", RED),
]

for i, (icon, name, desc, color) in enumerate(sources):
    x = Inches(1) + Inches(i * 1.95)
    y = Inches(2.8)
    card = add_shape(slide, x, y, Inches(1.75), Inches(2.2), DARK_BLUE)
    add_text(slide, x, y + Inches(0.2), Inches(1.75), Inches(0.5),
             icon, size=32, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.8), Inches(1.75), Inches(0.4),
             name, size=16, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.3), Inches(1.75), Inches(0.8),
             desc, size=12, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# Bottom note
add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
         "You choose which sources to scan. All data stays on your machine. Nothing is shared or uploaded.",
         size=16, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)

# Arrow pointing down
add_text(slide, Inches(5.5), Inches(6.2), Inches(2), Inches(0.5),
         "↓  YOUR CONTEXT.md  ↓", size=16, color=ACCENT_TEAL, bold=True,
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 6: Step 1 — Clone & Setup
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "STEP 1 OF 3", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Clone the Repo & Run Setup", size=36, color=WHITE, bold=True)

add_text(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.5),
         "Two commands. Under 2 minutes.", size=20, color=MEDIUM_GRAY)

# Code block background
add_shape(slide, Inches(1), Inches(3.2), Inches(11), Inches(2.5),
          RGBColor(0x0D, 0x11, 0x17))

code_lines = [
    "# Clone the PRISM",
    "git clone https://github.com/[YourGitHubUsername]/PRISM.git \\",
    "    ~/Documents/Claude/PRISM",
    "",
    "# Run the one-time setup",
    "cd ~/Documents/Claude/PRISM",
    "chmod +x setup.sh && ./setup.sh",
]

for i, line in enumerate(code_lines):
    color = MEDIUM_GRAY if line.startswith("#") else GREEN
    if not line:
        continue
    add_text(slide, Inches(1.4), Inches(3.4) + Inches(i * 0.32),
             Inches(10), Inches(0.3), line, size=16, color=color,
             font_name="Courier New")

# What happens
add_text(slide, Inches(1), Inches(6.0), Inches(10), Inches(0.4),
         "Setup creates your .personal/ folder, generates your CLAUDE.md, and sets up logging.",
         size=16, color=MEDIUM_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 7: Step 2 — Run Identity Scan
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "STEP 2 OF 3", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Run the Identity Scan", size=36, color=WHITE, bold=True)

add_text(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.5),
         "Open Claude Code and say one sentence.", size=20, color=MEDIUM_GRAY)

# Command block
add_shape(slide, Inches(1), Inches(3.2), Inches(11), Inches(1.2),
          RGBColor(0x0D, 0x11, 0x17))
add_text(slide, Inches(1.4), Inches(3.35), Inches(10), Inches(0.3),
         '> ', size=18, color=MEDIUM_GRAY, font_name="Courier New")
add_text(slide, Inches(1.8), Inches(3.35), Inches(10), Inches(0.3),
         '"Run the identity scan"', size=18, color=GREEN, font_name="Courier New")

add_text(slide, Inches(1.4), Inches(3.8), Inches(10), Inches(0.4),
         "Claude scans your environment and asks which sources to include.",
         size=14, color=MEDIUM_GRAY, font_name="Courier New")

# What gets generated
add_text(slide, Inches(1), Inches(4.8), Inches(10), Inches(0.4),
         "What gets generated:", size=18, color=WHITE, bold=True)

gen_items = [
    ("CONTEXT.md", "Your complete business profile — who you are, team, clients, priorities"),
    ("memory-bank/", "11 files covering identity, business, team, clients, relationships, tools"),
    ("config.yml", "Record of what was scanned and when — for future refreshes"),
]

for i, (name, desc) in enumerate(gen_items):
    y = Inches(5.3) + Inches(i * 0.55)
    add_text(slide, Inches(1.4), y, Inches(2.5), Inches(0.4),
             name, size=16, color=ACCENT_TEAL, bold=True, font_name="Courier New")
    add_text(slide, Inches(4), y, Inches(8), Inches(0.4),
             desc, size=16, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 8: Step 3 — Use It
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "STEP 3 OF 3", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Just Use Claude Code. The PRISM Does the Rest.",
         size=34, color=WHITE, bold=True)

# Flow diagram
steps = [
    ("You open Claude Code", MEDIUM_GRAY),
    ("→", ACCENT_TEAL),
    ("PRISM loads your context", ACCENT_BLUE),
    ("→", ACCENT_TEAL),
    ("You give a task", MEDIUM_GRAY),
    ("→", ACCENT_TEAL),
    ("SOP Router matches it", GREEN),
    ("→", ACCENT_TEAL),
    ("Claude follows the skill/SOP", ORANGE),
]

x = Inches(0.5)
for text, color in steps:
    if text == "→":
        add_text(slide, x, Inches(2.5), Inches(0.4), Inches(0.5),
                 text, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
        x += Inches(0.5)
    else:
        add_text(slide, x, Inches(2.5), Inches(1.8), Inches(0.7),
                 text, size=14, color=color, align=PP_ALIGN.CENTER)
        x += Inches(1.8)

# Example tasks
add_text(slide, Inches(1), Inches(3.8), Inches(10), Inches(0.4),
         "Example tasks the PRISM handles:", size=18, color=WHITE, bold=True)

examples = [
    ('"Write an article from this transcript"', '→  article-writer.md skill  →  18-step QA gate'),
    ('"Run Dollar-a-Day for this client"', '→  dollar-a-day.md skill  →  3-phase campaign setup'),
    ('"Build a KP sprint plan"', '→  kp-sprint.md skill  →  30-day timeline with milestones'),
    ('"Write the weekly MAA report"', '→  weekly-maa-report.md  →  Metrics, Analysis, Action format'),
    ('"Onboard a new client"', '→  client-onboarding SOP  →  GCT form + access checklist'),
]

for i, (task, result) in enumerate(examples):
    y = Inches(4.4) + Inches(i * 0.5)
    add_text(slide, Inches(1.4), y, Inches(4.5), Inches(0.4),
             task, size=14, color=ACCENT_TEAL, font_name="Courier New")
    add_text(slide, Inches(6), y, Inches(6.5), Inches(0.4),
             result, size=14, color=MEDIUM_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 9: What's Inside — Skills
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "WHAT'S INSIDE", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "12 Skills — Ready-to-Execute Action Prompts", size=34,
         color=WHITE, bold=True)

skills = [
    ("Content Factory", "6-stage content production pipeline"),
    ("Article Writer", "Transcript-to-article with voice matching"),
    ("Article QA", "18-step [Methodology Partner] quality gate"),
    ("Dollar-a-Day", "3-phase paid social campaign setup"),
    ("Content Repurposing", "Long-form to multi-platform distribution"),
    ("Weekly MAA Report", "Metrics, Analysis, Action framework"),
    ("KP Sprint", "30-day Knowledge Panel accelerator"),
    ("Influence Report Card", "Authority score and presence audit"),
    ("Personal Brand Site", "WordPress site build for KP clients"),
    ("Prospect Follow-up", "Discovery call prep and qualification"),
    ("SOP Router", "Auto-matches tasks to the right SOP"),
    ("PRISM Core", "System operating instructions"),
]

for i, (name, desc) in enumerate(skills):
    col = 0 if i < 6 else 1
    row = i if i < 6 else i - 6
    x = Inches(1) + Inches(col * 6)
    y = Inches(2.5) + Inches(row * 0.7)

    add_text(slide, x, y, Inches(0.3), Inches(0.4),
             "▸", size=14, color=ACCENT_BLUE)
    add_text(slide, x + Inches(0.3), y, Inches(2.2), Inches(0.4),
             name, size=16, color=WHITE, bold=True)
    add_text(slide, x + Inches(2.5), y, Inches(3.2), Inches(0.4),
             desc, size=14, color=MEDIUM_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 10: Canon Frameworks
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "THE CANON", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "15 Source-of-Truth Frameworks Built In", size=34,
         color=WHITE, bold=True)

add_text(slide, Inches(1), Inches(2.2), Inches(10), Inches(0.5),
         "If an SOP contradicts a canon file, the canon is correct. Always.",
         size=18, color=ORANGE)

# 9 Triangles visual
triangles = [
    ("WHY", [("SBP", "Specialist / Business / Partner"),
             ("ACC", "Awareness / Consideration / Conversion"),
             ("GCT", "Goals / Content / Targeting")]),
    ("HOW", [("MAA", "Metrics / Analysis / Action"),
             ("LDT", "Learn / Do / Teach"),
             ("CCS", "Content / Checklist / Software")]),
    ("WHAT", [("DDD", "Do / Delegate / Delete"),
              ("CID", "Communicate / Iterate / Delegate"),
              ("MOF", "Marketing / Operations / Finance")]),
]

tier_colors = {"WHY": ACCENT_BLUE, "HOW": ACCENT_TEAL, "WHAT": GREEN}

for col, (tier, tris) in enumerate(triangles):
    x = Inches(1) + Inches(col * 4)
    y = Inches(3.2)

    add_shape(slide, x, y, Inches(3.6), Inches(0.5), tier_colors[tier])
    add_text(slide, x, y + Inches(0.05), Inches(3.6), Inches(0.4),
             tier, size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    for i, (abbr, desc) in enumerate(tris):
        ty = y + Inches(0.7) + Inches(i * 0.7)
        add_text(slide, x + Inches(0.2), ty, Inches(0.8), Inches(0.4),
                 abbr, size=16, color=tier_colors[tier], bold=True)
        add_text(slide, x + Inches(1), ty, Inches(2.5), Inches(0.4),
                 desc, size=13, color=LIGHT_GRAY)

# Content Factory
add_text(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
         "Content Factory:  Plumbing → Produce → Process → Post → Promote → Perform",
         size=16, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(6.4), Inches(11), Inches(0.4),
         "Always 6 stages. Never 4 or 5.", size=14, color=MEDIUM_GRAY,
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 11: For Teams
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "FOR TEAMS", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Everyone Runs the Same System. Customized to Them.",
         size=34, color=WHITE, bold=True)

# Three person columns
people = [
    ("Founder / CEO", "Full business context, all clients,\nfinancial data, strategic priorities", ACCENT_BLUE),
    ("Content Specialist", "Assigned clients, content pipeline,\narticle QA standards, voice profiles", GREEN),
    ("Client / Partner", "Their business context only,\nrelevant SOPs, their team directory", ORANGE),
]

for i, (role, context, color) in enumerate(people):
    x = Inches(1) + Inches(i * 4)
    card = add_shape(slide, x, Inches(2.8), Inches(3.5), Inches(3), DARK_BLUE)

    # Icon circle
    icon_shape = add_shape(slide, x + Inches(1.3), Inches(3.0), Inches(0.8), Inches(0.8), color)

    add_text(slide, x, Inches(4.0), Inches(3.5), Inches(0.5),
             role, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(4.6), Inches(3), Inches(1),
             context, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Bottom note
add_text(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
         "Same SOPs. Same quality standards. Same methodology.\nDifferent personal context loaded on boot.",
         size=18, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 12: Compounding
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "COMPOUND, DON'T CAMPAIGN", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Every Session Makes the PRISM Smarter", size=36,
         color=WHITE, bold=True)

# Timeline
stages = [
    ("Day 1", "Identity scan builds\nyour initial context", "70%", ACCENT_BLUE),
    ("Week 1", "Memory bank fills in\nfrom daily work", "[Majority %]", ACCENT_TEAL),
    ("Week 2", "SOPs get updated with\nyour discoveries", "90%", GREEN),
    ("Month 1", "System knows you as well\nas a veteran assistant", "95%", ORANGE),
    ("Ongoing", "Auto-updates, never\nstarts from zero again", "99%", RGBColor(0xAB, 0x6E, 0xFF)),
]

for i, (time, desc, pct, color) in enumerate(stages):
    x = Inches(0.8) + Inches(i * 2.4)
    y = Inches(2.8)

    # Vertical line
    add_shape(slide, x + Inches(0.9), y, Inches(0.04), Inches(3.5), color)

    # Time label
    add_text(slide, x, y + Inches(0.2), Inches(2), Inches(0.4),
             time, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)

    # Description
    add_text(slide, x, y + Inches(0.8), Inches(2), Inches(0.8),
             desc, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Percentage
    add_text(slide, x, y + Inches(2.0), Inches(2), Inches(0.5),
             pct, size=28, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(2.6), Inches(2), Inches(0.4),
             "context accuracy", size=11, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════
# SLIDE 13: Privacy & Security
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "PRIVACY & SECURITY", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Your Data Stays on Your Machine. Period.", size=36,
         color=WHITE, bold=True)

privacy_items = [
    ("🔒", "Local Only", "All personal context lives in .personal/ on YOUR computer. Never uploaded."),
    ("🚫", "Git-Ignored", "Personal files are in .gitignore. They never touch the shared repo."),
    ("👁️", "You Choose", "The identity scan asks permission before scanning each source."),
    ("🗑️", "Deletable", "Delete .personal/ at any time to reset. Run the scan again to rebuild."),
    ("📖", "Read-Only Scan", "The identity scan only reads. It never modifies, sends, or publishes."),
]

for i, (icon, title, desc) in enumerate(privacy_items):
    y = Inches(2.6) + Inches(i * 0.9)
    add_text(slide, Inches(1.2), y, Inches(0.5), Inches(0.5),
             icon, size=24, color=WHITE)
    add_text(slide, Inches(2), y, Inches(2), Inches(0.4),
             title, size=18, color=WHITE, bold=True)
    add_text(slide, Inches(4.2), y, Inches(8), Inches(0.4),
             desc, size=16, color=LIGHT_GRAY)


# ════════════════════════════════════════════════════
# SLIDE 14: Quick Start Summary
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(0.5), Inches(10), Inches(0.6),
         "GET STARTED", size=14, color=ACCENT_TEAL, bold=True)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
         "Three Steps. Five Minutes. Permanently Better Claude.",
         size=36, color=WHITE, bold=True)

# Three big numbered steps
step_data = [
    ("1", "Clone & Setup", "git clone + ./setup.sh", "2 min", ACCENT_BLUE),
    ("2", "Identity Scan", '"Run the identity scan"', "5-15 min", ACCENT_TEAL),
    ("3", "Use It", "Every session is now\npersonalized", "∞", GREEN),
]

for i, (num, title, desc, time, color) in enumerate(step_data):
    x = Inches(1.5) + Inches(i * 3.8)
    y = Inches(2.8)

    # Number circle
    circle = add_shape(slide, x + Inches(0.8), y, Inches(1.2), Inches(1.2), color)

    add_text(slide, x + Inches(0.8), y + Inches(0.15), Inches(1.2), Inches(1),
             num, size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(1.5), Inches(2.8), Inches(0.5),
             title, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(2.1), Inches(2.8), Inches(0.6),
             desc, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    add_text(slide, x, y + Inches(2.9), Inches(2.8), Inches(0.4),
             time, size=14, color=color, align=PP_ALIGN.CENTER)

# Repo URL
add_shape(slide, Inches(3), Inches(6.2), Inches(7), Inches(0.7),
          RGBColor(0x0D, 0x11, 0x17))
add_text(slide, Inches(3), Inches(6.3), Inches(7), Inches(0.5),
         "github.com/[YourGitHubUsername]/PRISM", size=22, color=GREEN,
         bold=True, align=PP_ALIGN.CENTER, font_name="Courier New")


# ════════════════════════════════════════════════════
# SLIDE 15: Closing
# ════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLACK)
add_shape(slide, Inches(0), Inches(0), W, Inches(0.08), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.2),
         "Stop explaining yourself\nto Claude every session.", size=44,
         color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.6),
         "Let the PRISM do it for you.", size=28,
         color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_shape(slide, Inches(0), Inches(6.9), W, Inches(0.08), ACCENT_TEAL)
add_text(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.4),
         "[Your Agency Name]  ×  [Methodology Partner]  |  github.com/[YourGitHubUsername]/PRISM",
         size=12, color=MEDIUM_GRAY, align=PP_ALIGN.CENTER)


# ── Save ──
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "Claude-PRISM-Onboarding.pptx")
prs.save(output_path)
print(f"Deck saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
